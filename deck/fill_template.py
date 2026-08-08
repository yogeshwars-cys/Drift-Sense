"""Fill the official i4C Hackathon 2026 template with the Drift-Sense submission.

    python deck/fill_template.py

Reads  Idea-Submission-Template_Hackathon-2026-1.pptx  (unmodified, as issued)
Writes deck/DriftSense_Submission.pptx

Rules this follows, from the template's own instructions slide and from the
problem statement:

  * the instructions slide is deleted, as it says to
  * the template's own section headings are left exactly as issued -- only the
    {placeholder} bodies are filled, because the instructions forbid changing
    the idea-detail pointers
  * paragraphs are avoided in favour of points, as it asks
  * the problem statement separately requires ONE success and ONE honest
    failure visual, so the two cards on the Impact slide are replaced by those
    two figures; they are the only layout change made

Formatting is preserved by cloning the template's own first paragraph for each
line rather than assigning `text_frame.text`, which would collapse every run to
a single unstyled one.
"""
import copy
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
TEMPLATE = os.path.join(ROOT, 'Idea-Submission-Template_Hackathon-2026-1.pptx')
OUT = os.path.join(HERE, 'DriftSense_Submission.pptx')
REPO = 'https://github.com/yogeshwars-cys/Drift-Sense'


# --------------------------------------------------------------- helpers ----
def by_id(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f'shape {shape_id} not on this slide')


def set_lines(shape, lines, size=None):
    """Replace a text frame's content, keeping the template's own styling.

    The first paragraph is the style template: it is cloned once per line, so
    font, colour, bullet and spacing all survive. Assigning text_frame.text
    instead would flatten every run to an unstyled one."""
    tf = shape.text_frame
    proto = copy.deepcopy(tf.paragraphs[0]._p)
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]._p
    for r in first.findall(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}r')[1:]:
        first.remove(r)

    def fill(p_el, text):
        runs = p_el.findall(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
        if not runs:
            return
        for extra in runs[1:]:
            p_el.remove(extra)
        t = runs[0].find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
        t.text = text

    fill(first, lines[0])
    for line in lines[1:]:
        clone = copy.deepcopy(proto)
        fill(clone, line)
        first.getparent().append(clone)
    if size:
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size)


def place(slide, shape_id, x=None, y=None, w=None, h=None):
    """Move/resize a shape, in inches. The template sizes every body box for a
    one-line stub, so filling them with real content overflows unless the box
    and its surrounding card are grown into the panel's free space."""
    sh = by_id(slide, shape_id)
    if x is not None:
        sh.left = Inches(x)
    if y is not None:
        sh.top = Inches(y)
    if w is not None:
        sh.width = Inches(w)
    if h is not None:
        sh.height = Inches(h)


def drop(slide, *shape_ids):
    for sid in shape_ids:
        try:
            sh = by_id(slide, sid)
        except KeyError:
            continue
        sh._element.getparent().remove(sh._element)


# ------------------------------------------------------------------ main ----
def main():
    if not os.path.exists(TEMPLATE):
        sys.exit(f'template not found: {TEMPLATE}')
    prs = Presentation(TEMPLATE)

    # 1. structural first: delete the instructions slide, as it instructs
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])

    s = list(prs.slides)          # now: 0 team, 1 problem, 2 idea, 3 solution,
                                  # 4 innovation, 5 impact, 6 tech, 7 links, 8 refs

    # ------------------------------------------------------- 1. team ------
    # Rows 3 and 4 are dashed rather than left as {Enter Name}: the team is two
    # people, and a submitted deck should not carry unfilled template stubs.
    set_lines(by_id(s[0], 54), ['Silicon Bakers'])
    for sid, val in ((55, 'Yogeshwar S'), (56, '3rd Year'),
                     (57, 'Ankam Charam Teja'), (58, '3rd Year'),
                     (59, '—'), (60, '—'),
                     (61, '—'), (62, '—')):
        set_lines(by_id(s[0], sid), [val])
    # the name cells are 1.19in wide, sized for the {Enter Name} stub; the
    # column has clear run to the ACADEMIC YEAR column at x=9.29
    for sid in (55, 57, 59, 61):
        place(s[0], sid, w=4.40)
    set_lines(by_id(s[0], 63), ['Amrita Vishwa Vidyapeetham'])
    set_lines(by_id(s[0], 64), ['+91 94878 40396'])
    set_lines(by_id(s[0], 65), ['yogeshwarsetg@gmail.com'])

    # ---------------------------------------------------- 2. problem ------
    place(s[1], 16, w=11.2, h=0.5)
    place(s[1], 17, y=4.15, h=2.80)          # card
    place(s[1], 19, y=4.72, w=10.8, h=2.15)  # body
    set_lines(by_id(s[1], 16),
              ['Selected: Drift-Sense - AI-Powered Navigation-Error Recovery'])
    set_lines(by_id(s[1], 19), [
        'A tool must revisit the same site thousands of times a day. Between '
        'visits the stage drifts - thermal, vibration, mechanical slack - and '
        'lands several pixels off.',
        'It cannot detect this from the image: every die repeats the same '
        'layout, so the wrong cell looks identical to the right one.',
        'Template matching fails exactly where layouts are most regular - ~30 '
        'correlation peaks within 0.003 of the maximum in a single window.',
        'Consequence: measurements trended over time silently compare two '
        'different places on the die.',
    ], size=13)

    # ------------------------------------------------------- 3. idea ------
    place(s[2], 17, y=4.30, h=1.34)
    place(s[2], 18, y=4.52)
    place(s[2], 19, y=4.82, w=10.8, h=0.78)
    place(s[2], 20, y=5.74, h=1.30)
    place(s[2], 21, y=5.96)
    place(s[2], 22, y=6.26, w=10.8, h=0.70)
    set_lines(by_id(s[2], 19), [
        'Measure the geometry, do not search it: the ~10x magnification IS the '
        'ratio of the two lattice pitches - recovered to 0.10% median.',
        'Spectrally subtract the lattice; what survives (array edges, dropped '
        'vias, gate crossings) is the only content carrying absolute identity.',
    ], size=12)
    set_lines(by_id(s[2], 22), [
        'Propose cheaply over the frame, then decide at FULL reference '
        'resolution - where cells actually differ. Candidates compete.',
        'DRAM and FinFET, 50/50. Classical CV, zero learned weights on the '
        'shipped path - a measured choice, not an omission.',
    ], size=12)

    # --------------------------------------------------- 4. solution ------
    place(s[3], 17, y=4.30, h=2.65)
    place(s[3], 18, y=4.55)
    place(s[3], 19, y=4.88, w=10.8, h=2.00)
    set_lines(by_id(s[3], 19), [
        'GENERATOR - DRAM grids and FinFET fin/gate fields rendered '
        'supersampled, then area-averaged down 10x.',
        'Imaging chain in physical order: edge brightening, astigmatic PSF, '
        'scan distortion, charging, shading, sensor noise last.',
        'Independent Poisson+Gaussian noise per capture, search side heavier. '
        'Rotation on the reference crop only - rotating the search canvas '
        'silently corrupts ground truth.',
        'Ground truth records the true centre AND whether the site is '
        'identifiable at all.',
        'PIPELINE - lattice sensor > sub-bin rotation + phase lock > aperiodic '
        'residual > NCC proposal + point-source votes > full-res rescoring > '
        'decision > sub-pixel refinement > (x, y).',
    ], size=12)

    # ------------------------------------------------- 5. innovation ------
    for cid, bid in ((17, 19), (20, 22)):
        place(s[4], cid, y=4.21, h=2.74)
        place(s[4], bid, y=4.76, w=4.9, h=2.10)
    set_lines(by_id(s[4], 19), [
        'Magnification is MEASURED, not swept - a blind 8-point sweep becomes '
        'a tight bracket.',
        'Spatial induction: the lattice must prove its own geometry. Catches '
        '13/13 silent failures at zero false alarms - and, sign inverted, '
        'predicts solvability (84.6% vs a 50% base rate).',
        'Point-source voting: a dropped via is 1.3% of a footprint-wide '
        'correlation and is never proposed. A blob detector votes for the '
        'centre instead - recovering those sites with zero losses.',
    ], size=11)
    set_lines(by_id(s[4], 22), [
        'Template matching slides a template and takes the argmax. We measure '
        'the geometry, remove the lattice, and let candidates compete.',
        'The system knows what it cannot know: 100% recall on detecting sites '
        'that are information-theoretically unidentifiable.',
        'Every idea measured that did NOT pay is documented, not deleted - the '
        'centre prior, the phase penalty, the commit gate, a CNN re-ranker, a '
        'learned ranker. Five nulls, with the evidence shown.',
    ], size=11)

    # ------------------------------------- 6. impact / results + visuals --
    place(s[5], 16, y=3.24, w=11.2, h=1.05)
    set_lines(by_id(s[5], 16), [
        'ACCURACY  50.0% of 100 pairs within 15px; 62.5% of the 80 identifiable '
        'pairs; 92.5% when a non-periodic landmark is in view.',
        'SPEED  3.1s per 1000x1000 pair on ONE CPU core - no GPU, no weights.  '
        'GEOMETRY  magnification 0.10%, rotation 0.15 deg.',
        'ROBUSTNESS  tripling search-image sensor noise costs nothing '
        'measurable.',
    ], size=12)
    drop(s[5], 17, 18, 19, 20, 21, 22, 23, 24)
    figs = [('success_case.png', 'SUCCESS - array corner in view, error 0.4px'),
            ('failure_case.png', 'HONEST FAILURE - defect-free array interior, '
                                 'error 320.5px at confidence 0.85: wrong AND sure')]
    for i, (fn, cap) in enumerate(figs):
        path = os.path.join(ROOT, 'examples', fn)
        if not os.path.exists(path):
            continue
        x = Inches(1.05 + i * 5.75)
        s[5].shapes.add_picture(path, x, Inches(4.42),
                                width=Inches(5.40), height=Inches(2.42))
        tb = s[5].shapes.add_textbox(x, Inches(6.86), Inches(5.40), Inches(0.30))
        tf = tb.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = cap
        r.font.size = Pt(8)
        r.font.bold = True

    # -------------------------------------------------------- 7. tech ----
    place(s[6], 17, y=3.44, h=2.16)
    place(s[6], 18, y=3.66)
    place(s[6], 19, y=3.96, w=9.9, h=1.55)
    set_lines(by_id(s[6], 19), [
        'STACK  Python 3.12 - NumPy, OpenCV, SciPy, Pillow. Four packages, no '
        'GPU, no torch on the inference path, no weights to download.',
        'HARDWARE  every result produced on a single consumer CPU core.',
        'TIMING  generation ~200s per 100 pairs; inference 3.1s mean, 7.1s '
        'worst case. Shipped model size: 0 bytes.',
        'FEASIBILITY  two commands: pip install -r requirements-inference.txt, '
        'then python localize.py ref.png search.png.',
        'REPRODUCIBILITY  requirements.txt is the complete pip freeze; every '
        'figure regenerates from a recorded seed.',
    ], size=11)

    # ------------------------------------------------------- 8. links ----
    set_lines(by_id(s[7], 20), [REPO])
    set_lines(by_id(s[7], 27),
              ['{Paste your video link here - localize.py on one sample pair}'])

    # -------------------------------------------------------- 9. refs ----
    place(s[8], 18, y=3.04, h=1.30)
    place(s[8], 19, y=3.22, w=10.1, h=1.05)
    set_lines(by_id(s[8], 19), [
        'Every augmentation, noise model and structural parameter is justified '
        'in CITATIONS.md - 15 sections, 2-3 public references each.',
        'Research foundation: a perfectly periodic field is translation-'
        'invariant modulo the pitch, so absolute position is not recoverable '
        'from pixels. Measured - array corner in view 92.5%, nothing aperiodic '
        '0.0%.',
    ], size=11)
    place(s[8], 24, y=4.98, h=1.85)
    for sid, yy in ((26, 5.12), (28, 5.68), (29, 6.24)):
        place(s[8], sid, y=yy, w=10.1, h=0.50)
    refs = [
        'Zhang et al. (2019), A Poisson-Gaussian Denoising Dataset with Real '
        'Fluorescence Microscopy Images, CVPR - arxiv.org/pdf/1812.10366',
        'US Patent 9,430,457 - Ambiguity Reduction for Image Alignment '
        'Applications: NCC over a repeating pattern yields indistinguishable '
        'peaks, requiring explicit disambiguation.',
        'US Patent 11,481,922 - Online Navigational Drift Correction for '
        'Metrology Measurements: the same thermal-drift navigation problem.',
    ]
    for sid, txt in zip((26, 28, 29), refs):
        set_lines(by_id(s[8], sid), [txt], size=10)
    prs.save(OUT)
    print(f'wrote {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)')


if __name__ == '__main__':
    main()
